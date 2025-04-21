import pandas as pd
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches

# Load the Excel file
excel_file = '../../Budget.xlsx'
sheet_name = 'Charts'

# Function to extract tables from a sheet
def extract_tables(df):
    tables = []
    start_row = None
    for i, row in df.iterrows():
        if row.isnull().all():
            if start_row is not None:
                tables.append(df.iloc[start_row:i])
                start_row = None
        else:
            if start_row is None:
                start_row = i
    if start_row is not None:
        tables.append(df.iloc[start_row:])
    return tables

# Read the "Charts" sheet into a DataFrame
df = pd.read_excel(excel_file, sheet_name=sheet_name)

# Extract tables from the sheet
tables = extract_tables(df)

# Create a PowerPoint presentation
prs = Presentation()

for idx, table in enumerate(tables):
    # Print the content of the table for verification
    print(f"Content of table {idx+1} in sheet {sheet_name}:")
    print(table)
    
    # Skip the initial rows to get to the actual data
    data_start_row = 1  # Adjust this based on your data structure
    
    # Extract X values (first row of actual data, next to unnamed columns)
    x_values = table.iloc[data_start_row, 2:15].astype(str)  # Convert X values to strings
    
    # Extract Y names (second column, starting from the row after the heading)
    y_names = table.iloc[data_start_row+1:, 1]  # Y names from the second column, starting from the row after the heading
    
    # Extract Y values (from the row after the heading onwards, starting from the third column up to the 15th column)
    y_values = table.iloc[data_start_row+1:, 2:15]
    
    # Plot each Y series
    plt.figure(figsize=(10, 6))
    for i, y_name in enumerate(y_names):
        y_series = y_values.iloc[i]
        plt.plot(x_values, y_series, label=y_name)
        # Annotate each data point with its value
        for x, y in zip(x_values, y_series):
            plt.text(x, y, f'{y:.2f}', fontsize=8, ha='right')
    
    # Add labels and title
    plt.xlabel('X Values')
    plt.ylabel('Y Values')
    plt.title(f'{table.iloc[data_start_row-1, 0]} for {sheet_name}')
    plt.legend()
    
    # Save the plot as an image
    chart_image = os.path.abspath(f'{sheet_name}_chart_{idx+1}.png')
    plt.savefig(chart_image)
    plt.close()

    print(f'Chart saved as {chart_image}')
    
    # Add a slide to the PowerPoint presentation
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the slide title
    title = slide.shapes.title
    title.text = str(table.iloc[data_start_row-1, 0])
    
    # Add the chart image to the slide
    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(chart_image, left, top, width=Inches(8), height=Inches(4.5))

# Save the PowerPoint presentation
ppt_file = 'charts_presentation.pptx'
prs.save(ppt_file)

print(f'Presentation saved as {ppt_file}')